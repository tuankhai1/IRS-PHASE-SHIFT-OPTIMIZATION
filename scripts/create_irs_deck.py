import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / ".codex_deps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "outputs"
OUT_FILE = OUT_DIR / "irs_component_optimization_slides.pptx"
RESULTS = ROOT / "results"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

BLACK = RGBColor(18, 18, 18)
GRAY = RGBColor(92, 96, 104)
LIGHT = RGBColor(242, 244, 247)
PANEL = RGBColor(237, 239, 242)
RULE = RGBColor(184, 188, 196)
BLUE = RGBColor(45, 125, 246)
CYAN = RGBColor(0, 172, 193)
GREEN = RGBColor(67, 160, 71)
ORANGE = RGBColor(251, 140, 0)
PURPLE = RGBColor(142, 36, 170)
RED = RGBColor(229, 57, 53)


def set_text(tf, text, size=18, color=BLACK, bold=False, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for r in p.runs:
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def add_text(slide, text, left, top, width, height, size=18, color=BLACK,
             bold=False, align=None, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    set_text(tf, text, size, color, bold, align)
    if line_spacing:
        tf.paragraphs[0].line_spacing = line_spacing
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, Inches(0.55), Inches(0.33), Inches(11.8),
             Inches(0.55), size=35, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55),
                                  Inches(1.05), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    if subtitle:
        add_text(slide, subtitle, Inches(0.58), Inches(1.15), Inches(12),
                 Inches(0.35), size=16, color=GRAY)


def add_footer(slide, n):
    add_text(slide, f"{n:02d}", Inches(12.55), Inches(7.04), Inches(0.35),
             Inches(0.2), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height, fill=PANEL, line=RGBColor(222, 225, 230)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_bullets(slide, items, left, top, width, height, size=18, color=BLACK,
                bullet_color=BLUE):
    y = top
    for item in items:
        add_text(slide, "•", left, y + Inches(0.02), Inches(0.22),
                 Inches(0.25), size=size, color=bullet_color, bold=True)
        add_text(slide, item, left + Inches(0.27), y, width - Inches(0.27),
                 Inches(0.42), size=size, color=color)
        y += Inches(0.5)


def add_metric(slide, label, value, left, top, width, accent=BLUE):
    add_panel(slide, left, top, width, Inches(1.05), fill=RGBColor(248, 249, 251))
    add_text(slide, value, left + Inches(0.18), top + Inches(0.12),
             width - Inches(0.36), Inches(0.36), size=24, bold=True, color=accent)
    add_text(slide, label, left + Inches(0.18), top + Inches(0.57),
             width - Inches(0.36), Inches(0.3), size=13, color=GRAY)


def add_image_contain(slide, path, left, top, width, height):
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = width / height
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        w = width
        h = width / img_ratio
        x = left
        y = top + (height - h) / 2
    else:
        h = height
        w = height * img_ratio
        x = left + (width - w) / 2
        y = top
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_table(slide, rows, left, top, width, height, font_size=14):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    tbl = table_shape.table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Arial"
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = BLACK if i else RGBColor(255, 255, 255)
                    r.font.bold = i == 0
            fill = BLUE if i == 0 else (RGBColor(248, 249, 251) if i % 2 else RGBColor(255, 255, 255))
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
    return table_shape


def add_result_slide(prs, n, title, fig_name, claim, bullets, metrics=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_title(slide, title, claim)
    add_image_contain(slide, RESULTS / fig_name, Inches(0.65), Inches(1.55),
                      Inches(8.25), Inches(5.15))
    add_panel(slide, Inches(9.25), Inches(1.55), Inches(3.35), Inches(3.05),
              fill=RGBColor(248, 249, 251))
    add_text(slide, "Ý nghĩa chính", Inches(9.45), Inches(1.78),
             Inches(2.9), Inches(0.32), size=20, bold=True)
    add_bullets(slide, bullets, Inches(9.45), Inches(2.25),
                Inches(2.85), Inches(1.7), size=16)
    if metrics:
        x = Inches(9.25)
        y = Inches(4.85)
        for label, value, color in metrics:
            add_metric(slide, label, value, x, y, Inches(1.55), color)
            x += Inches(1.75)
    add_footer(slide, n)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    # 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_text(s, "Tối ưu tham số linh kiện IRS", Inches(0.65), Inches(1.65),
             Inches(11.5), Inches(0.75), size=50, bold=True)
    add_text(s, "Tối đa hóa tốc độ truyền dữ liệu với mô hình phản xạ thực tế",
             Inches(0.68), Inches(2.55), Inches(10.5), Inches(0.45),
             size=24, color=GRAY)
    add_panel(s, Inches(0.7), Inches(4.2), Inches(11.85), Inches(1.25),
              fill=RGBColor(248, 249, 251))
    add_text(s, "Chuỗi tối ưu:  L1,n, L2,n, Cn, Rn  →  Zn  →  vn  →  RSE",
             Inches(1.0), Inches(4.58), Inches(11.1), Inches(0.42),
             size=25, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_footer(s, 1)

    # 2
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Bài toán không chỉ là chọn pha phản xạ",
              "Tối ưu linh kiện cho phép mô hình hóa đúng ràng buộc phần cứng của IRS.")
    add_text(s, "Câu hỏi trung tâm", Inches(0.75), Inches(1.65),
             Inches(3.5), Inches(0.35), size=22, bold=True)
    add_text(s, "Làm thế nào chọn linh kiện của từng phần tử IRS để cực đại tốc độ truyền dữ liệu?",
             Inches(0.75), Inches(2.15), Inches(4.55), Inches(1.25),
             size=26, bold=True, color=BLACK)
    add_panel(s, Inches(5.75), Inches(1.55), Inches(6.85), Inches(4.75),
              fill=RGBColor(248, 249, 251))
    add_bullets(s, [
        "Mô hình lý tưởng giả sử |vn| = 1 nên chỉ cần chỉnh pha.",
        "Mô hình thực tế có suy hao biên độ phụ thuộc pha.",
        "Tối ưu trực tiếp linh kiện tạo nghiệm khả thi với mạch IRS.",
        "Đánh giá bằng tốc độ truyền đạt được RSE."
    ], Inches(6.15), Inches(2.0), Inches(6.0), Inches(2.7), size=20)
    add_footer(s, 2)

    # 3
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Mỗi nghiệm ứng viên là một vector 4N biến",
              "Thuật toán tìm kiếm trong không gian linh kiện thay vì không gian pha.")
    add_panel(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(1.0),
              fill=RGBColor(248, 249, 251))
    add_text(s, "x = [L1,1, L2,1, C1, R1, ..., L1,N, L2,N, CN, RN]",
             Inches(1.0), Inches(1.9), Inches(11.4), Inches(0.36),
             size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    rows = [
        ["Biến", "Ý nghĩa", "Miền giá trị"],
        ["L1,n", "Coupling inductance", "0.5 - 5.0 nH"],
        ["L2,n", "Varactor series inductance", "0.1 - 3.0 nH"],
        ["Cn", "Varactor capacitance", "0.1 - 5.0 pF"],
        ["Rn", "Series resistance", "0.5 - 5.0 Ω"],
    ]
    add_table(s, rows, Inches(0.9), Inches(3.05), Inches(7.05),
              Inches(2.65), font_size=15)
    add_panel(s, Inches(8.35), Inches(3.05), Inches(3.75), Inches(2.65),
              fill=RGBColor(248, 249, 251))
    add_text(s, "Ràng buộc", Inches(8.6), Inches(3.32), Inches(3.2),
             Inches(0.32), size=21, bold=True)
    add_text(s, "Tất cả biến được chặn trong miền vật lý. Khi thuật toán cập nhật vượt biên, nghiệm được đưa về giá trị gần nhất trong miền cho phép.",
             Inches(8.6), Inches(3.85), Inches(3.2), Inches(1.15),
             size=17, color=BLACK)
    add_text(s, "f = 5.8 GHz,  ω = 2πf,  Z0 = 377 Ω",
             Inches(8.6), Inches(5.15), Inches(3.1), Inches(0.35),
             size=16, color=GRAY)
    add_footer(s, 3)

    # 4
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Từ linh kiện sang hệ số phản xạ của phần tử IRS",
              "Trở kháng Zn là cầu nối giữa thiết kế mạch và tín hiệu phản xạ.")
    add_panel(s, Inches(0.75), Inches(1.7), Inches(5.7), Inches(2.05),
              fill=RGBColor(248, 249, 251))
    add_text(s, "Zn = jωL1,n (jωL2,n + 1/(jωCn) + Rn)\n     / (jωL1,n + jωL2,n + 1/(jωCn) + Rn)",
             Inches(1.05), Inches(2.16), Inches(5.1), Inches(0.85),
             size=21, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_panel(s, Inches(6.95), Inches(1.7), Inches(5.6), Inches(2.05),
              fill=RGBColor(248, 249, 251))
    add_text(s, "vn = (Zn - Z0) / (Zn + Z0)\n|vn| = βn,     angle(vn) = θn",
             Inches(7.25), Inches(2.2), Inches(5.0), Inches(0.75),
             size=23, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_bullets(s, [
        "L2,n và Cn điều chỉnh phần phản kháng nên ảnh hưởng mạnh tới pha.",
        "Rn biểu diễn tổn hao, làm giảm biên độ phản xạ.",
        "Một thay đổi linh kiện có thể làm đổi đồng thời βn và θn."
    ], Inches(1.05), Inches(4.25), Inches(11.2), Inches(1.8), size=20)
    add_footer(s, 4)

    # 5
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Fitness của thuật toán là tốc độ truyền đạt được",
              "Mỗi nghiệm linh kiện được chuyển thành vn trước khi tính RSE.")
    add_panel(s, Inches(0.75), Inches(1.6), Inches(11.85), Inches(1.25),
              fill=RGBColor(248, 249, 251))
    add_text(s, "RSE = log2(1 + PT ||vᴴΦ + hdᴴ||² / σ²)",
             Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.42),
             size=27, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "Bài toán tối ưu", Inches(0.95), Inches(3.45),
             Inches(3.0), Inches(0.32), size=22, bold=True)
    add_text(s, "maximize_x  RSE(x)\nsubject to physical bounds of L1,n, L2,n, Cn, Rn",
             Inches(0.95), Inches(3.95), Inches(5.25), Inches(0.9),
             size=21, color=BLACK)
    add_text(s, "Các đường so sánh", Inches(6.9), Inches(3.45),
             Inches(3.6), Inches(0.32), size=22, bold=True)
    add_bullets(s, [
        "Upper bound: IRS lý tưởng.",
        "AO practical: tối ưu pha theo mô hình tương đương.",
        "Component-level: tối ưu trực tiếp L1, L2, C, R.",
        "Lower bound: không dùng IRS."
    ], Inches(6.9), Inches(3.95), Inches(5.5), Inches(1.8), size=18)
    add_footer(s, 5)

    # 6
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "GWO component-level là lựa chọn chính cho chất lượng nghiệm",
              "PSO nhẹ hơn, nhưng GWO đạt tốc độ truyền cao hơn trong các kết quả mô phỏng.")
    add_panel(s, Inches(0.75), Inches(1.6), Inches(3.75), Inches(4.6),
              fill=RGBColor(248, 249, 251))
    add_text(s, "Vì sao dùng GWO?", Inches(1.0), Inches(1.95),
             Inches(3.2), Inches(0.35), size=21, bold=True)
    add_bullets(s, [
        "Không cần gradient.",
        "Phù hợp bài toán phi lồi.",
        "Xử lý trực tiếp ràng buộc hộp.",
        "Tìm kiếm tốt trong không gian 4N biến."
    ], Inches(1.0), Inches(2.55), Inches(3.1), Inches(2.1), size=17)
    add_panel(s, Inches(4.85), Inches(1.6), Inches(3.75), Inches(4.6),
              fill=RGBColor(248, 249, 251))
    add_text(s, "Cơ chế GWO", Inches(5.1), Inches(1.95),
             Inches(3.2), Inches(0.35), size=21, bold=True)
    add_text(s, "Ba nghiệm tốt nhất α, β, δ dẫn hướng quần thể.\n\nTham số a giảm từ 2 về 0: đầu quá trình ưu tiên exploration, cuối quá trình ưu tiên exploitation.",
             Inches(5.1), Inches(2.55), Inches(3.1), Inches(2.45),
             size=18)
    add_panel(s, Inches(8.95), Inches(1.6), Inches(3.65), Inches(4.6),
              fill=RGBColor(248, 249, 251))
    add_text(s, "Kết luận chọn thuật toán", Inches(9.2), Inches(1.95),
             Inches(3.0), Inches(0.35), size=21, bold=True)
    add_text(s, "Dùng GWO khi cần hiệu năng cao nhất; dùng PSO khi muốn giảm thời gian chạy. Hybrid PSO-GWO có thể cải thiện thêm ở một số cấu hình.",
             Inches(9.2), Inches(2.55), Inches(3.0), Inches(1.8),
             size=18)
    add_footer(s, 6)

    # 7-14 result slides
    add_result_slide(
        prs, 7, "Fig. 5: AO thực tế bám sát xu hướng upper bound theo khoảng cách",
        "fig5_rate_vs_distance.png",
        "Tốc độ tăng mạnh khi user tiến gần IRS; mô hình lý tưởng vẫn là trần hiệu năng.",
        [
            "AO practical cao hơn lower bound rõ rệt.",
            "Ideal design khi đánh giá thực tế bị mất hiệu năng.",
            "Tại d = 500 m, upper bound đạt khoảng 4.55 bit/s/Hz."
        ],
        [("AO practical", "3.46", GREEN), ("No IRS", "0.15", GRAY)]
    )
    add_result_slide(
        prs, 8, "Fig. 6: Tăng số phần tử IRS giúp tăng tốc độ truyền",
        "fig6_rate_vs_N.png",
        "N lớn hơn tạo thêm bậc tự do phản xạ, nhưng vẫn chịu giới hạn phần cứng thực tế.",
        [
            "Upper bound tăng từ 1.01 lên 5.01 bit/s/Hz.",
            "AO practical tăng từ 0.66 lên 3.84 bit/s/Hz.",
            "Lower bound gần như không đổi vì không dùng IRS."
        ],
        [("N = 80 AO", "3.84", GREEN), ("N = 80 ideal", "5.01", RED)]
    )
    add_result_slide(
        prs, 9, "Fig. 7: Lượng tử pha càng mịn thì hiệu năng càng tốt",
        "fig7_discrete_phases.png",
        "Rời rạc pha làm giảm hiệu năng; tăng số bit pha giúp tiến gần mô hình liên tục hơn.",
        [
            "b = 3 tốt hơn b = 1 và b = 2.",
            "Khoảng cách với upper bound vẫn tồn tại do suy hao thực tế.",
            "Hiệu ứng rõ nhất khi user gần IRS."
        ],
        [("Practical b=3", "2.21", GREEN), ("Ideal b=3", "2.92", BLUE)]
    )
    add_result_slide(
        prs, 10, "Fig. 8: Tối ưu linh kiện có thể vượt AO tối ưu pha thực tế",
        "fig8_component_vs_distance.png",
        "GWO component-level khai thác miền linh kiện tốt hơn khi khoảng cách thuận lợi.",
        [
            "GWO component gần upper bound hơn PSO.",
            "PSO component vẫn vượt AO tại các điểm d lớn.",
            "Tại d = 500 m, GWO đạt khoảng 4.00 bit/s/Hz."
        ],
        [("GWO comp.", "4.00", CYAN), ("AO practical", "3.42", GREEN)]
    )
    add_result_slide(
        prs, 11, "Fig. 9: Khi N tăng, lợi thế của tối ưu linh kiện rõ hơn",
        "fig9_component_vs_N.png",
        "Không gian 4N biến lớn hơn, nhưng cũng cho phép điều chỉnh phản xạ linh hoạt hơn.",
        [
            "GWO component đạt 4.13 bit/s/Hz tại N = 80.",
            "PSO component đạt 3.72 bit/s/Hz tại N = 80.",
            "AO practical đạt 3.82 bit/s/Hz tại N = 80."
        ],
        [("GWO N=80", "4.13", CYAN), ("Upper", "5.00", RED)]
    )
    add_result_slide(
        prs, 12, "Fig. 10: GWO hội tụ tới nghiệm tốt hơn PSO/APSO",
        "fig10_convergence.png",
        "Đường hội tụ cho thấy GWO tiếp tục cải thiện sau khi PSO gần bão hòa.",
        [
            "PSO tăng nhanh nhưng chững lại sớm hơn.",
            "APSO không vượt PSO trong cấu hình này.",
            "GWO đạt giá trị cuối cao nhất."
        ],
        [("GWO final", "2.81", CYAN), ("PSO final", "2.48", PURPLE)]
    )
    add_result_slide(
        prs, 13, "Fig. 11: Hybrid PSO-GWO cải thiện nhẹ so với GWO thuần",
        "fig11_phase_vs_component.png",
        "Kết hợp khởi tạo/tìm kiếm giúp tiến gần upper bound hơn ở nhiều điểm khoảng cách.",
        [
            "Hybrid PSO-GWO đạt 4.06 bit/s/Hz tại d = 500 m.",
            "GWO component đạt 4.00 bit/s/Hz.",
            "PSO component thấp hơn nhưng chạy nhanh hơn."
        ],
        [("Hybrid", "4.06", GREEN), ("GWO", "4.00", CYAN)]
    )
    add_result_slide(
        prs, 14, "Fig. 12: Cn là tham số then chốt trong tối ưu linh kiện",
        "fig12_fixed_component_ablation.png",
        "Cố định C và R làm tốc độ giảm mạnh; chỉ cố định R gây mất mát nhỏ hơn.",
        [
            "Full optimization đạt 3.55 bit/s/Hz tại d = 500 m.",
            "Fix R chỉ giảm xuống 3.46 bit/s/Hz.",
            "Fix C,R giảm còn 1.71 bit/s/Hz."
        ],
        [("Full", "3.55", PURPLE), ("Fix C,R", "1.71", GREEN)]
    )

    # 15
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Phân tích ảnh hưởng linh kiện đến tốc độ truyền",
              "Kết quả ablation cho thấy không phải linh kiện nào cũng quan trọng như nhau.")
    rows = [
        ["Cấu hình tại d = 500 m", "RSE (bit/s/Hz)", "Diễn giải"],
        ["Tối ưu L1,L2,C,R", "3.55", "Mở đủ bậc tự do của mạch"],
        ["Cố định R", "3.46", "Mất mát nhỏ, R chủ yếu gây suy hao"],
        ["Cố định C,R", "1.71", "Mất mạnh khả năng điều chỉnh cộng hưởng"],
        ["Chỉ tối ưu L1", "1.35", "Không đủ tự do điều chỉnh vn"],
        ["Không IRS", "0.15", "Mốc dưới"],
    ]
    add_table(s, rows, Inches(0.75), Inches(1.7), Inches(7.45),
              Inches(3.8), font_size=13)
    add_panel(s, Inches(8.6), Inches(1.7), Inches(3.95), Inches(3.8),
              fill=RGBColor(248, 249, 251))
    add_text(s, "Thông điệp", Inches(8.9), Inches(2.05),
             Inches(3.3), Inches(0.32), size=22, bold=True)
    add_text(s, "Cn ảnh hưởng trực tiếp tới phần dung kháng 1/(jωCn), nên quyết định mạnh tới cộng hưởng, pha phản xạ và hiệu năng tổng thể.\n\nRn quan trọng cho suy hao, nhưng trong miền khảo sát việc cố định R gây giảm nhẹ hơn.",
             Inches(8.9), Inches(2.6), Inches(3.25), Inches(2.0),
             size=17)
    add_footer(s, 15)

    # 16
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Kết luận: tối ưu linh kiện làm rõ giới hạn phần cứng IRS",
              "Cách tiếp cận component-level liên kết trực tiếp thiết kế mạch với tốc độ truyền.")
    add_bullets(s, [
        "Biểu diễn nghiệm bằng vector [L1,n, L2,n, Cn, Rn] giúp nghiệm luôn gắn với linh kiện khả thi.",
        "Chuỗi tính Zn → vn → RSE cho phép đánh giá đúng ảnh hưởng đồng thời của biên độ và pha.",
        "GWO component-level đạt hiệu năng cao hơn PSO trong các kết quả chính, đổi lại thời gian chạy lớn hơn.",
        "So với AO tối ưu pha theo mô hình tương đương, tối ưu linh kiện cho thấy rõ vai trò của Cn và các ràng buộc vật lý.",
        "Khi một số linh kiện bị cố định, hiệu năng giảm mạnh nhất nếu mất khả năng điều chỉnh Cn."
    ], Inches(1.0), Inches(1.75), Inches(11.4), Inches(3.2), size=21)
    add_panel(s, Inches(1.0), Inches(5.55), Inches(11.4), Inches(0.65),
              fill=RGBColor(248, 249, 251))
    add_text(s, "Kết quả mô phỏng sử dụng 300 channel realizations cho mỗi điểm trên trục x.",
             Inches(1.25), Inches(5.75), Inches(10.9), Inches(0.25),
             size=16, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(s, 16)

    OUT_DIR.mkdir(exist_ok=True)
    prs.save(OUT_FILE)
    return OUT_FILE


if __name__ == "__main__":
    path = build_deck()
    print(path)
